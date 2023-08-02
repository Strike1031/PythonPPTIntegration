import sys
import os
import copy
import pptx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.shapes.picture import Picture
from pptx.enum.shapes import MSO_SHAPE

#                                                               #
#       Main Function to insert a Slide to another file         #       
#                                                               #

def func(tempfile, workfile, pageNumber, result_pptx = "result.pptx"):
    tempPPT = Presentation(tempfile)
    workPPT = Presentation(workfile)
    tempSlide = tempPPT.slides[0]
    newSlide = workPPT.slides.add_slide(tempSlide.slide_layout)

    
    for m in newSlide.shapes:
        if m.is_placeholder:
            sp = m.placeholder_format
            if sp.idx == 0:  
                m.text = "" 
            elif sp.idx == 1: 
                m.text = "" 
            else:
                newSlide.shapes._spTree.remove(m._element)


    images = {} #Image List
    newSlide.slide_layout.description = ""

    
    #Save files in the template file
    for m in tempSlide.shapes:
        if 'Picture' in m.name:
            with open(m.name+'.jpg', 'wb') as f:
                f.write(m.image.blob)
            
            images[m.name+'.jpg'] = [m.left, m.top, m.width, m.height]
        else:
            el = m.element
            newel = copy.deepcopy(el)
            newSlide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    # Integrate pictures into the target file
    for i, j in images.items():
        image = newSlide.shapes.add_picture(i, j[0], j[1], j[2], j[3])
        image.width = j[2]
        image.height = j[3]
        os.remove(i) 

    slide_id = workPPT.slides.index(newSlide)
 
    xml_slides = workPPT.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[slide_id])
    xml_slides.insert(int(pageNumber), slides[slide_id])  
   

    workPPT.save(result_pptx)
    print("Success!!!")

def main():
    tempfile = "temp.pptx"  # temp file
    workfile = "work.pptx"  # working file
    pageNumber = 2  # Page to Insert
    result_pptx = "result.pptx"  # Result
    tempfile = input("Enter template file name: ") # temp file
    workfile = input("Enter working file name: ") # working file
    pageNumber = int(input("Page Number to insert: ")) # Page to Insert
    result_pptx = input("Enter output file name: ") # Result
    func(tempfile, workfile, pageNumber, result_pptx)

if __name__ == '__main__':
    main()